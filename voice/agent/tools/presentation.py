"""
agent/tools/presentation.py
Create or update PowerPoint presentations using content from
another document (PDF, text, markdown, existing slides).
"""
from __future__ import annotations

import json
import logging
import os
import re
from pathlib import Path

from agent import llm as llm_cfg
from agent.tools.document import (
    DOCUMENTS_DIR,
    HOME,
    MAX_CHARS,
    _extract_json_object,
    read_pdf,
)

log = logging.getLogger("intellivox.presentation")


def _read_plain_text(path: str) -> dict:
    expanded = str(Path(path).expanduser().resolve())
    if not os.path.exists(expanded):
        return {"success": False, "message": f"File not found: {path}"}
    try:
        with open(expanded, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        return {
            "success": True,
            "text": text[:MAX_CHARS],
            "path": expanded,
            "truncated": len(text) > MAX_CHARS,
        }
    except Exception as e:
        return {"success": False, "message": f"Could not read file: {e}"}


def read_document_text(path: str) -> dict:
    """
    Read text from PDF, plain text/markdown, or PowerPoint.
    Returns: { success, text, path, kind }
    """
    expanded = str(Path(path).expanduser().resolve())
    if not os.path.exists(expanded):
        return {"success": False, "message": f"File not found: {path}"}

    ext = Path(expanded).suffix.lower()
    if ext == ".pdf":
        result = read_pdf(expanded)
        if result.get("success"):
            result["kind"] = "pdf"
        return result

    if ext in (".pptx", ".ppt"):
        outline = _presentation_outline(expanded)
        if not outline.get("success"):
            return outline
        return {
            "success": True,
            "text": outline["outline_text"],
            "path": expanded,
            "kind": "presentation",
            "slide_count": outline.get("slide_count", 0),
        }

    if ext in (".txt", ".md", ".markdown", ".csv", ".json", ".log", ""):
        result = _read_plain_text(expanded)
        if result.get("success"):
            result["kind"] = "text"
        return result

    # Best-effort plain read for unknown types
    result = _read_plain_text(expanded)
    if result.get("success") and (result.get("text") or "").strip():
        result["kind"] = "text"
        return result
    return {
        "success": False,
        "message": f"Unsupported document type for text extraction: {ext or '(none)'}",
    }


def _presentation_outline(path: str) -> dict:
    """Extract title + body text from each slide."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER
    except ImportError:
        return {
            "success": False,
            "message": "python-pptx is required. Run: pip install python-pptx",
        }

    expanded = str(Path(path).expanduser().resolve())
    try:
        prs = Presentation(expanded)
    except Exception as e:
        return {"success": False, "message": f"Could not open presentation: {e}"}

    slides = []
    lines = []
    for i, slide in enumerate(prs.slides):
        title = ""
        bullets: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text_frame.text or "").strip()
            if not text:
                continue
            is_title = False
            if shape.is_placeholder:
                try:
                    ph = shape.placeholder_format.type
                    is_title = ph in (
                        PP_PLACEHOLDER.TITLE,
                        PP_PLACEHOLDER.CENTER_TITLE,
                        PP_PLACEHOLDER.VERTICAL_TITLE,
                    )
                except Exception:
                    is_title = False
            if is_title and not title:
                title = text.splitlines()[0].strip()
            else:
                for line in text.splitlines():
                    line = line.strip().lstrip("•-*–— ").strip()
                    if line:
                        bullets.append(line)
        if not title and bullets:
            title = bullets[0]
            bullets = bullets[1:]
        slides.append({"index": i, "title": title or f"Slide {i + 1}", "bullets": bullets[:12]})
        lines.append(f"Slide {i + 1}: {slides[-1]['title']}")
        for b in bullets[:8]:
            lines.append(f"  - {b}")

    return {
        "success": True,
        "path": expanded,
        "slides": slides,
        "slide_count": len(slides),
        "outline_text": "\n".join(lines) if lines else "(empty presentation)",
    }


def _normalize_slides(data) -> list[dict]:
    """Normalize LLM JSON into a list of {title, bullets}."""
    if data is None:
        return []

    raw_slides = None
    if isinstance(data, dict):
        raw_slides = data.get("slides") or data.get("deck") or data.get("presentation")
        if raw_slides is None and ("title" in data or "bullets" in data):
            raw_slides = [data]
    elif isinstance(data, list):
        raw_slides = data

    if not isinstance(raw_slides, list):
        return []

    out: list[dict] = []
    for item in raw_slides:
        if isinstance(item, str):
            out.append({"title": item.strip()[:80] or "Slide", "bullets": []})
            continue
        if not isinstance(item, dict):
            continue
        title = str(item.get("title") or item.get("heading") or "Slide").strip()[:120]
        bullets = item.get("bullets") or item.get("points") or item.get("content") or []
        if isinstance(bullets, str):
            bullets = [ln.strip("•-* ") for ln in bullets.splitlines() if ln.strip()]
        elif isinstance(bullets, list):
            flat = []
            for b in bullets:
                if isinstance(b, dict):
                    flat.append(str(b.get("text") or b.get("point") or next(iter(b.values()), "")).strip())
                else:
                    flat.append(str(b).strip())
            bullets = [b for b in flat if b]
        else:
            bullets = [str(bullets).strip()] if bullets else []
        out.append({"title": title or "Slide", "bullets": bullets[:10]})
    return out


def _fallback_slides_from_text(text: str, max_slides: int = 6) -> list[dict]:
    """Heuristic slide split when the LLM returns unusable JSON."""
    chunks = [c.strip() for c in re.split(r"\n\s*\n", text) if c.strip()]
    if not chunks:
        # Fall back to line-based chunks
        lines = [ln.strip("•-*–— ") for ln in text.splitlines() if ln.strip()]
        if not lines:
            return [{"title": "Overview", "bullets": ["No content extracted"]}]
        # Group every ~4 lines into a slide
        chunks = []
        for i in range(0, len(lines), 4):
            chunks.append("\n".join(lines[i : i + 4]))

    slides = []
    for i, chunk in enumerate(chunks[:max_slides]):
        lines = [ln.strip("•-*–— ") for ln in chunk.splitlines() if ln.strip()]
        title = (lines[0][:80] if lines else f"Slide {i + 1}")
        bullets = lines[1:6] if len(lines) > 1 else lines[:5]
        slides.append({"title": title, "bullets": bullets})
    return slides or [{"title": "Overview", "bullets": [text[:200]]}]


def _collect_source_paths(
    source_path: str = "",
    document_path: str = "",
    source_paths=None,
    document_paths=None,
) -> list[str]:
    """Normalize one or many source document paths into a unique list."""
    paths: list[str] = []

    def _add(value) -> None:
        if value is None:
            return
        if isinstance(value, (list, tuple)):
            for item in value:
                _add(item)
            return
        text = str(value).strip().strip("\"'")
        if not text:
            return
        # Allow comma / semicolon / newline separated lists in a single string
        candidate = str(Path(text).expanduser())
        if any(sep in text for sep in (",", ";", "\n")) and not os.path.exists(candidate):
            for part in re.split(r"[,;\n]+", text):
                _add(part)
            return
        if text not in paths:
            paths.append(text)

    _add(source_paths)
    _add(document_paths)
    _add(source_path)
    _add(document_path)
    return paths


def _read_sources(paths: list[str]) -> dict:
    """Read one or more documents and merge their text."""
    if not paths:
        return {"success": False, "message": "No source document path provided"}

    texts: list[str] = []
    resolved: list[str] = []
    kinds: list[str] = []
    errors: list[str] = []

    for path in paths:
        result = read_document_text(path)
        if not result.get("success"):
            errors.append(result.get("message") or f"Failed to read {path}")
            continue
        body = (result.get("text") or "").strip()
        if not body:
            errors.append(f"No extractable text in {result.get('path', path)}")
            continue
        label = Path(result["path"]).name
        texts.append(f"--- Source: {label} ---\n{body}")
        resolved.append(result["path"])
        kinds.append(result.get("kind") or "text")

    if not texts:
        return {
            "success": False,
            "message": "; ".join(errors) if errors else "Source documents have no extractable text",
        }

    merged = "\n\n".join(texts)
    return {
        "success": True,
        "text": merged[:MAX_CHARS],
        "path": resolved[0],
        "paths": resolved,
        "kinds": kinds,
        "truncated": len(merged) > MAX_CHARS,
        "message": f"Read {len(resolved)} source document(s)",
    }


def _write_pptx(slides: list[dict], path: str) -> dict:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {
            "success": False,
            "message": "python-pptx is required. Run: pip install python-pptx",
        }

    expanded = str(Path(path).expanduser().resolve())
    if not expanded.startswith(HOME):
        return {"success": False, "message": "Access denied: outside home directory"}
    Path(expanded).parent.mkdir(parents=True, exist_ok=True)
    if not expanded.lower().endswith(".pptx"):
        expanded = str(Path(expanded).with_suffix(".pptx"))

    prs = Presentation()
    # Use blank widescreen-ish default; title+content layout index 1 is standard
    layout_title = prs.slide_layouts[0]  # title slide
    layout_content = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    for i, slide_data in enumerate(slides):
        layout = layout_title if i == 0 else layout_content
        slide = prs.slides.add_slide(layout)
        title = slide_data.get("title") or f"Slide {i + 1}"
        bullets = slide_data.get("bullets") or []

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Body placeholder
        body = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body = shape
                break
        if body is None:
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    body = shape
                    break

        if body is not None and bullets:
            tf = body.text_frame
            tf.clear()
            for bi, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0
                try:
                    p.font.size = Pt(20)
                except Exception:
                    pass
        elif bullets and body is None:
            # Add a text box if no body placeholder
            box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.5), Inches(4.5))
            tf = box.text_frame
            for bi, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                try:
                    p.font.size = Pt(20)
                except Exception:
                    pass

    try:
        prs.save(expanded)
    except Exception as e:
        return {"success": False, "message": f"Could not save presentation: {e}"}

    return {
        "success": True,
        "path": expanded,
        "slide_count": len(slides),
        "message": f"Saved {len(slides)} slides to {expanded}",
    }


def _plan_slides_with_llm(
    source_text: str,
    current_outline: str = "",
    query: str = "",
    max_slides: int = 8,
) -> tuple[list[dict], str]:
    """Ask the LLM for an updated slide deck. Returns (slides, raw)."""
    task = (query or "").strip() or (
        "Update the presentation so it reflects the latest facts from the source document(s). "
        "Keep useful structure from the current deck when present, but replace outdated content."
    )
    current_block = current_outline.strip() or "(no existing slides — create a new deck)"

    prompt = (
        "You update a slide deck using information from source document(s).\n"
        "Return ONLY valid JSON in this exact shape:\n"
        '{"slides":[{"title":"Slide title","bullets":["point 1","point 2"]}]}\n'
        "Rules:\n"
        f"- Produce between 3 and {max_slides} slides (never only 1 slide if there is enough content)\n"
        "- First slide should be a title/overview slide\n"
        "- Each later slide covers one theme from the sources\n"
        "- Bullets must be short (max ~12 words each), max 5 bullets per slide\n"
        "- Prefer facts from the SOURCE document(s) over the old deck\n"
        "- Do not invent facts that are not supported by the source\n"
        f"Task: {task}\n\n"
        f"=== CURRENT PRESENTATION ===\n{current_block[:MAX_CHARS // 3]}\n\n"
        f"=== SOURCE DOCUMENT(S) ===\n{source_text[: (2 * MAX_CHARS) // 3]}"
    )

    raw = ""
    try:
        raw = llm_cfg.chat(
            [{"role": "user", "content": prompt}],
            temperature=0.2,
            num_predict=900,
            format="json",
        )
    except Exception as e:
        log.warning("JSON-format chat failed (%s); retrying plain", e)
        try:
            raw = llm_cfg.chat(
                [{"role": "user", "content": prompt}],
                temperature=0.2,
                num_predict=900,
            )
        except Exception as e2:
            log.error("Slide planning LLM failed: %s", e2)
            return _fallback_slides_from_text(source_text, max_slides=max_slides), ""

    slides = _normalize_slides(_extract_json_object(raw))
    # Small models often return a single mega-slide — expand with heuristic if needed
    if not slides or (len(slides) < 3 and len(source_text) > 200):
        fallback = _fallback_slides_from_text(source_text or raw, max_slides=max_slides)
        if not slides:
            slides = fallback
        elif len(fallback) > len(slides):
            # Keep LLM titles when useful, but prefer richer fallback coverage
            slides = fallback
    return slides[:max_slides], raw


def update_presentation_from_document(
    source_path: str = "",
    presentation_path: str = "",
    query: str = "",
    output_path: str = "",
    open_after: bool = True,
    document_path: str = "",
    ppt_path: str = "",
    source_paths=None,
    document_paths=None,
) -> dict:
    """
    Update (or create) a PowerPoint presentation using information from other document(s).

    source_path / document_path: PDF/text/etc. providing the new information
    source_paths / document_paths: optional list (or comma-separated) of extra sources
    presentation_path / ppt_path: existing .pptx to update (optional — creates new if missing)
    query: optional focus (e.g. "emphasize pricing and deadlines")
    output_path: where to save; default ~/Documents/<name>_updated.pptx

    Returns: {
      success, path, slide_count, slides, source_path, source_paths, presentation_path,
      opened, message
    }
    """
    paths = _collect_source_paths(
        source_path=source_path,
        document_path=document_path,
        source_paths=source_paths,
        document_paths=document_paths,
    )
    presentation_path = (presentation_path or ppt_path or "").strip()

    if not paths:
        return {"success": False, "message": "No source document path provided"}

    source = _read_sources(paths)
    if not source.get("success"):
        return source

    source_text = (source.get("text") or "").strip()
    if not source_text:
        return {"success": False, "message": "Source document has no extractable text"}

    current_outline = ""
    existing_path = ""
    if presentation_path:
        expanded_ppt = str(Path(presentation_path).expanduser().resolve())
        if os.path.exists(expanded_ppt):
            existing_path = expanded_ppt
            outline = _presentation_outline(expanded_ppt)
            if outline.get("success"):
                current_outline = outline.get("outline_text", "")
            else:
                log.warning("Could not read existing deck: %s", outline.get("message"))
        else:
            log.info("Presentation not found at %s — will create a new deck", presentation_path)

    slides, _raw = _plan_slides_with_llm(
        source_text=source_text,
        current_outline=current_outline,
        query=query,
    )
    if not slides:
        return {"success": False, "message": "Could not build slides from the source document"}

    out = (output_path or "").strip()
    if not out:
        if existing_path:
            stem = Path(existing_path).stem
            out = str(Path(existing_path).with_name(f"{stem}_updated.pptx"))
        else:
            src_stem = Path(source["path"]).stem
            safe = re.sub(r"[^\w\-]+", "_", src_stem).strip("_") or "presentation"
            DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)
            out = str(DOCUMENTS_DIR / f"{safe}_presentation.pptx")

    written = _write_pptx(slides, out)
    if not written.get("success"):
        return written

    opened = False
    open_msg = ""
    if open_after:
        try:
            from agent import platform as plat

            open_res = plat.open_path(written["path"])
            opened = bool(open_res.get("success"))
            open_msg = open_res.get("message", "")
        except Exception as e:
            open_msg = str(e)

    action = "Updated" if existing_path else "Created"
    src_label = ", ".join(Path(p).name for p in source.get("paths", [source["path"]]))
    msg = (
        f"{action} presentation ({written['slide_count']} slides) from {src_label} "
        f"→ {written['path']}"
    )
    if opened:
        msg += " (opened)"
    elif open_after and open_msg:
        msg += f" — open note: {open_msg}"

    log.info("%s presentation at %s from %s", action, written["path"], source.get("paths"))
    return {
        "success": True,
        "path": written["path"],
        "slide_count": written["slide_count"],
        "slides": slides,
        "source_path": source["path"],
        "source_paths": source.get("paths", [source["path"]]),
        "presentation_path": existing_path or "",
        "opened": opened,
        "message": msg,
    }


def create_presentation_from_document(
    source_path: str = "",
    query: str = "",
    output_path: str = "",
    open_after: bool = True,
    document_path: str = "",
    source_paths=None,
    document_paths=None,
) -> dict:
    """Create a new presentation from one or more source documents (no existing deck)."""
    return update_presentation_from_document(
        source_path=source_path or document_path,
        source_paths=source_paths,
        document_paths=document_paths,
        presentation_path="",
        query=query or "Create a clear presentation summarizing the key points",
        output_path=output_path,
        open_after=open_after,
    )
